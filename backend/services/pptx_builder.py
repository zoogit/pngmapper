from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from services.coordinate_converter import MapCoordinateConverter
from services.standard_map import get_standard_map_path, get_map_bounds, ASPECT_RATIOS, generate_map, REGION_BOUNDS
import os

# Alaska and Hawaii bounds for insets
ALASKA_BOUNDS = {
    'north': 71.5, 'south': 51.0,
    'west': -168.0, 'east': -130.0
}

HAWAII_BOUNDS = {
    'north': 22.5, 'south': 18.5,
    'west': -160.5, 'east': -154.5
}

def separate_us_locations(locations):
    """
    Separate locations into continental US, Alaska, and Hawaii groups

    Args:
        locations: List of dicts with 'lat' and 'lng'

    Returns:
        dict with keys 'continental', 'alaska', 'hawaii'
    """
    continental = []
    alaska = []
    hawaii = []

    for loc in locations:
        lat = loc.get('lat')
        lng = loc.get('lng')

        # Alaska: north of 51°N and west of 130°W
        if lat and lng and lat > 51.0 and lng < -130.0:
            alaska.append(loc)
        # Hawaii: south of 22°N and west of 155°W
        elif lat and lng and lat < 22.0 and lng < -155.0:
            hawaii.append(loc)
        # Continental US
        else:
            continental.append(loc)

    return {
        'continental': continental,
        'alaska': alaska,
        'hawaii': hawaii
    }

def create_presentation(map_image_path, locations):
    """
    Create PowerPoint presentation with map

    Args:
        map_image_path: Path to map image
        locations: List of location dicts

    Returns:
        str: Path to created presentation
    """
    prs = Presentation()

    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "Location Map"
    subtitle.text = f"{len(locations)} locations"

    # Map Slide
    blank_slide_layout = prs.slide_layouts[6]
    map_slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(12.333)
    height = Inches(0.6)

    title_box = map_slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Location Overview"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Add map image
    img_left = Inches(0.5)
    img_top = Inches(1.2)
    img_width = Inches(12.333)

    map_slide.shapes.add_picture(
        map_image_path,
        img_left,
        img_top,
        width=img_width
    )

    # Locations List Slide (if there are named locations)
    named_locations = [loc for loc in locations if loc.get('name')]
    if named_locations:
        list_slide = prs.slides.add_slide(prs.slide_layouts[1])
        list_title = list_slide.shapes.title
        list_title.text = "Locations List"

        body_shape = list_slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        for i, loc in enumerate(named_locations[:10], 1):  # Limit to first 10
            p = text_frame.add_paragraph()
            p.text = f"{i}. {loc['name']}"
            p.level = 0

    # Save
    output_path = 'output.pptx'
    prs.save(output_path)

    return output_path


def create_presentation_with_shapes(location_sets=None, locations=None, template_path=None, map_bounds=None, marker_styles=None,
                                   region='us', aspect_ratio='widescreen', projection='web_mercator'):
    """
    Create PowerPoint presentation with shapes instead of images

    Args:
        location_sets: List of location set dicts, each with 'name', 'locations', and 'markerStyles'
        locations: (Legacy) List of location dicts with lat, lng, name
        template_path: Optional path to template PPTX with map background
        map_bounds: Optional dict with custom map bounds
        marker_styles: (Legacy) Optional dict with marker styling options
        region: Region code ('us', 'europe', 'world', etc.)
        aspect_ratio: 'widescreen' (16:9) or 'standard' (4:3)
        projection: Projection type ('web_mercator', 'robinson', 'equal_earth')

    Returns:
        str: Path to created presentation
    """
    # Handle both old single-set format and new multi-set format
    if location_sets is None:
        # Legacy format - convert to new format
        location_sets = [{
            'name': 'Set 1',
            'locations': locations or [],
            'markerStyles': marker_styles or {}
        }]

    # Flatten all locations for map bounds calculation
    all_locations = []
    for loc_set in location_sets:
        all_locations.extend(loc_set['locations'])
    # Get aspect ratio dimensions
    aspect = ASPECT_RATIOS.get(aspect_ratio, ASPECT_RATIOS['widescreen'])

    # SLIDE 1: User's template (if exists)
    if template_path and os.path.exists(template_path):
        # Load the user's template as base
        prs = Presentation(template_path)
        print(f"DEBUG: Loaded template from {template_path}")

        # For template slide, we need to match the map positioning
        # Try to detect if there's an image on the slide and use its bounds
        # Otherwise, assume the map uses the same aspect ratio fitting as slide 2
        if len(prs.slides) > 0:
            template_slide = prs.slides[0]

            # Get map bounds for template slide
            print(f"DEBUG: Getting map bounds for template slide...")
            template_map_path, template_bounds = get_standard_map_path(
                region=region,
                aspect_ratio=aspect_ratio,
                projection=projection,
                locations=all_locations
            )
            print(f"DEBUG: Template bounds: N={template_bounds['north']:.2f}, S={template_bounds['south']:.2f}")

            # Calculate letterboxing for template (match slide 2)
            # Get natural map dimensions
            from PIL import Image
            with Image.open(template_map_path) as img:
                map_aspect = img.width / img.height

            template_slide_width = aspect['width']
            template_slide_height = aspect['height']

            # Calculate letterboxed dimensions (same as slide 2)
            template_img_width = template_slide_width
            template_img_height = template_img_width / map_aspect

            if template_img_height > template_slide_height:
                template_img_height = template_slide_height
                template_img_width = template_img_height * map_aspect

            # Center the map
            template_img_left = (template_slide_width - template_img_width) / 2
            template_img_top = (template_slide_height - template_img_height) / 2

            print(f"DEBUG: Slide 1 - Map positioned at: left={template_img_left:.2f}\", top={template_img_top:.2f}\"")
            print(f"DEBUG: Slide 1 - Map size: {template_img_width:.2f}\" x {template_img_height:.2f}\"")

            # Create converter for template slide - uses letterboxed positioning
            template_converter = MapCoordinateConverter(
                map_bounds=template_bounds,  # Original bounds, no adjustment
                slide_bounds={
                    'left': template_img_left,
                    'top': template_img_top,
                    'width': template_img_width,
                    'height': template_img_height
                },
                projection=projection
            )

            # Add markers for each location set
            for loc_set in location_sets:
                add_markers_to_slide(
                    template_slide,
                    loc_set['locations'],
                    template_converter,
                    loc_set['markerStyles']
                )
    else:
        # No user template, create new presentation
        prs = Presentation()
        prs.slide_width = Inches(aspect['width'])
        prs.slide_height = Inches(aspect['height'])
        print(f"DEBUG: Created new presentation with {aspect_ratio} aspect ratio")

    # SLIDE 2: Generated map (always add this)
    # For US region, check if we need Alaska/Hawaii insets
    use_insets = False
    continental_locs = []
    alaska_locs = []
    hawaii_locs = []

    if region == 'us':
        # Separate locations by region
        separated = separate_us_locations(all_locations)
        continental_locs = separated['continental']
        alaska_locs = separated['alaska']
        hawaii_locs = separated['hawaii']

        # Use insets if we have Alaska or Hawaii locations
        use_insets = len(alaska_locs) > 0 or len(hawaii_locs) > 0

        if use_insets:
            print(f"DEBUG: Using inset approach - Continental: {len(continental_locs)}, Alaska: {len(alaska_locs)}, Hawaii: {len(hawaii_locs)}")
            # Force continental bounds for main map
            continental_bounds = REGION_BOUNDS['us']['continental']
            map_bounds = generate_map(
                bounds=continental_bounds,
                projection=projection,
                output_path='continental_map.png'
            )
            standard_map_path = 'continental_map.png'
        else:
            # No Alaska/Hawaii - use standard approach
            standard_map_path, map_bounds = get_standard_map_path(
                region=region,
                aspect_ratio=aspect_ratio,
                projection=projection,
                locations=all_locations
            )
    else:
        # Non-US regions - use standard approach
        standard_map_path, map_bounds = get_standard_map_path(
            region=region,
            aspect_ratio=aspect_ratio,
            projection=projection,
            locations=all_locations
        )

    # Create a blank slide for the generated map
    blank_layout = prs.slide_layouts[6] if prs.slide_layouts else None
    if blank_layout:
        map_slide_2 = prs.slides.add_slide(blank_layout)
    else:
        # Fallback: just add slide without layout
        map_slide_2 = prs.slides.add_slide(prs.slide_layouts[0])

    # Get slide dimensions
    slide_width = aspect['width']
    slide_height = aspect['height']

    # Add main map with letterboxing (maintain aspect ratio)
    padding = 0
    max_width = Inches(slide_width - (2 * padding))

    # Add the image (PowerPoint will maintain aspect ratio)
    pic = map_slide_2.shapes.add_picture(
        standard_map_path,
        Inches(0),  # Temp position
        Inches(0),  # Temp position
        width=max_width
    )

    # Get actual dimensions after PowerPoint scaled it
    actual_width = pic.width.inches
    actual_height = pic.height.inches

    # Check if height fits on slide
    if actual_height > slide_height - (2 * padding):
        # Too tall - resize based on height instead
        # Calculate aspect ratio and set both dimensions
        aspect_ratio_img = actual_width / actual_height
        actual_height = slide_height - (2 * padding)
        actual_width = actual_height * aspect_ratio_img

        pic.width = Inches(actual_width)
        pic.height = Inches(actual_height)

    # Center the image on the slide
    img_left = (slide_width - actual_width) / 2
    img_top = (slide_height - actual_height) / 2

    pic.left = Inches(img_left)
    pic.top = Inches(img_top)

    print(f"DEBUG: Slide 2 - Main map positioned at: left={img_left:.2f}\", top={img_top:.2f}\"")
    print(f"DEBUG: Slide 2 - Main map size: {actual_width:.2f}\" x {actual_height:.2f}\"")
    print(f"DEBUG: Using bounds: N={map_bounds['north']:.2f}, S={map_bounds['south']:.2f}")

    # Converter uses ORIGINAL geographic bounds and EXACT image position
    main_converter = MapCoordinateConverter(
        map_bounds=map_bounds,  # Original bounds, no adjustment
        slide_bounds={
            'left': img_left,
            'top': img_top,
            'width': actual_width,
            'height': actual_height
        },
        projection=projection
    )

    # Add markers to main map
    if use_insets:
        # Only add continental markers to main map
        for loc_set in location_sets:
            # Filter to only continental locations
            continental_from_set = [loc for loc in loc_set['locations']
                                   if loc in continental_locs]
            if continental_from_set:
                add_markers_to_slide(
                    map_slide_2,
                    continental_from_set,
                    main_converter,
                    loc_set['markerStyles']
                )
    else:
        # Add all markers to main map (standard behavior)
        for loc_set in location_sets:
            add_markers_to_slide(
                map_slide_2,
                loc_set['locations'],
                main_converter,
                loc_set['markerStyles']
            )

    # Add Alaska inset if needed
    if use_insets and len(alaska_locs) > 0:
        print(f"DEBUG: Adding Alaska inset with {len(alaska_locs)} locations")

        # Generate Alaska map
        generate_map(
            bounds=ALASKA_BOUNDS,
            projection=projection,
            output_path='alaska_inset.png'
        )
        alaska_map_path = 'alaska_inset.png'

        # Inset size: 20% of main map width
        inset_width = actual_width * 0.20

        # Add Alaska inset (top-left corner with padding)
        alaska_pic = map_slide_2.shapes.add_picture(
            alaska_map_path,
            Inches(img_left + 0.2),  # Padding from left
            Inches(img_top + 0.2),  # Padding from top
            width=Inches(inset_width)
        )

        # Already positioned at top-left
        inset_height = alaska_pic.height.inches

        print(f"DEBUG: Alaska inset - size: {inset_width:.2f}\" x {inset_height:.2f}\"")

        # Create converter for Alaska inset
        alaska_converter = MapCoordinateConverter(
            map_bounds=ALASKA_BOUNDS,
            slide_bounds={
                'left': alaska_pic.left.inches,
                'top': alaska_pic.top.inches,
                'width': alaska_pic.width.inches,
                'height': alaska_pic.height.inches
            },
            projection=projection
        )

        # Add Alaska markers
        for loc_set in location_sets:
            alaska_from_set = [loc for loc in loc_set['locations']
                              if loc in alaska_locs]
            if alaska_from_set:
                add_markers_to_slide(
                    map_slide_2,
                    alaska_from_set,
                    alaska_converter,
                    loc_set['markerStyles']
                )

    # Add Hawaii inset if needed
    if use_insets and len(hawaii_locs) > 0:
        print(f"DEBUG: Adding Hawaii inset with {len(hawaii_locs)} locations")

        # Generate Hawaii map
        generate_map(
            bounds=HAWAII_BOUNDS,
            projection=projection,
            output_path='hawaii_inset.png'
        )
        hawaii_map_path = 'hawaii_inset.png'

        # Inset size: 20% of main map width
        inset_width = actual_width * 0.20

        # Add Hawaii inset (bottom-right corner with padding)
        hawaii_pic = map_slide_2.shapes.add_picture(
            hawaii_map_path,
            Inches(0),  # Temp position
            Inches(0),  # Temp position
            width=Inches(inset_width)
        )

        # Position at bottom-right
        inset_height = hawaii_pic.height.inches
        hawaii_pic.left = Inches(img_left + actual_width - inset_width - 0.2)  # Padding from right
        hawaii_pic.top = Inches(img_top + actual_height - inset_height - 0.2)  # Padding from bottom

        print(f"DEBUG: Hawaii inset - size: {inset_width:.2f}\" x {inset_height:.2f}\"")

        # Create converter for Hawaii inset
        hawaii_converter = MapCoordinateConverter(
            map_bounds=HAWAII_BOUNDS,
            slide_bounds={
                'left': hawaii_pic.left.inches,
                'top': hawaii_pic.top.inches,
                'width': hawaii_pic.width.inches,
                'height': hawaii_pic.height.inches
            },
            projection=projection
        )

        # Add Hawaii markers
        for loc_set in location_sets:
            hawaii_from_set = [loc for loc in loc_set['locations']
                              if loc in hawaii_locs]
            if hawaii_from_set:
                add_markers_to_slide(
                    map_slide_2,
                    hawaii_from_set,
                    hawaii_converter,
                    loc_set['markerStyles']
                )

    # Save
    output_path = 'output.pptx'
    prs.save(output_path)
    print(f"DEBUG: Saved presentation to {output_path}")

    return output_path


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def add_markers_to_slide(slide, locations, converter, marker_styles=None):
    """
    Add location markers to a slide

    Args:
        slide: PowerPoint slide object
        locations: List of location dicts
        converter: MapCoordinateConverter instance
        marker_styles: Optional dict with marker styling options
    """
    # Default styles
    default_styles = {
        'markerColor': '#dc3545',
        'markerShape': 'circle',
        'markerSize': 0.1,
        'showFill': True,
        'outlineColor': '#ffffff',
        'outlineWidth': 1.0,
        'showOutline': True,
        'showShadow': False,
        'showLabels': False,
        'labelFontSize': 10,
        'labelTextColor': '#000000',
        'labelBold': True
    }

    # Use provided styles or defaults
    styles = marker_styles if marker_styles else default_styles

    # Convert colors from hex to RGB
    marker_rgb = hex_to_rgb(styles.get('markerColor', default_styles['markerColor']))
    outline_rgb = hex_to_rgb(styles.get('outlineColor', default_styles['outlineColor']))
    label_text_rgb = hex_to_rgb(styles.get('labelTextColor', default_styles['labelTextColor']))

    # Add location markers as shapes
    for location in locations:
        lat = location.get('lat')
        lng = location.get('lng')
        name = location.get('name', '')

        if lat is None or lng is None:
            continue

        # Convert lat/lng to slide position
        left, top = converter.lat_lng_to_slide(lat, lng)

        # Determine shape type
        marker_shape = styles.get('markerShape', default_styles['markerShape'])
        shape_map = {
            'circle': MSO_SHAPE.OVAL,
            'square': MSO_SHAPE.RECTANGLE,
            'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
            'star': MSO_SHAPE.STAR_5_POINT
        }
        shape_type = shape_map.get(marker_shape, MSO_SHAPE.OVAL)

        # Add shape
        marker_size = Inches(styles.get('markerSize', default_styles['markerSize']))
        shape = slide.shapes.add_shape(
            shape_type,
            left - marker_size/2,  # Center the shape
            top - marker_size/2,
            marker_size,
            marker_size
        )

        # Style the marker
        if styles.get('showFill', default_styles['showFill']):
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*marker_rgb)
        else:
            shape.fill.background()

        if styles.get('showOutline', default_styles['showOutline']):
            shape.line.color.rgb = RGBColor(*outline_rgb)
            shape.line.width = Pt(styles.get('outlineWidth', default_styles['outlineWidth']))
        else:
            shape.line.fill.background()

        # Handle shadow
        if not styles.get('showShadow', default_styles['showShadow']):
            shape.shadow.inherit = False
            shape.shadow.visible = False

        # Add label if name exists and labels are enabled
        if name and styles.get('showLabels', default_styles['showLabels']):
            label_width = Inches(2)
            label_height = Inches(0.4)
            label_box = slide.shapes.add_textbox(
                left + marker_size/2 + Inches(0.1),
                top - label_height/2,
                label_width,
                label_height
            )
            text_frame = label_box.text_frame
            text_frame.text = name
            text_frame.margin_bottom = Pt(0)
            text_frame.margin_top = Pt(0)
            text_frame.margin_left = Pt(5)
            text_frame.margin_right = Pt(5)

            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(styles.get('labelFontSize', default_styles['labelFontSize']))
            paragraph.font.bold = styles.get('labelBold', default_styles['labelBold'])
            paragraph.font.color.rgb = RGBColor(*label_text_rgb)

            # No background fill or outline for labels
            label_box.fill.background()
            label_box.line.fill.background()
